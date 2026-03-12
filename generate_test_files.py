import pandas as pd
from docx import Document
from pptx import Presentation
# Try fpdf for basic PDF, if not, reportlab
import os

try:
    from fpdf import FPDF
except ImportError:
    os.system("pip install fpdf")
    from fpdf import FPDF

# 1. CSV
csv_data = {
    "Project": ["Center Lane", "Stargate"],
    "Scientist": ["Karl Pribram", "Itzhak Bentov"],
    "Concept": ["Holographic Universe", "Gateway Process"]
}
df = pd.DataFrame(csv_data)
df.to_csv("test.csv", index=False)
print("Created test.csv")

# 2. XLSX
df.to_excel("test.xlsx", index=False)
print("Created test.xlsx")

# 3. DOCX
doc = Document()
doc.add_heading('Test Document: McDonnell Report Snippet', 0)
doc.add_paragraph('This document contains details about Project Center Lane.')
doc.add_paragraph('Key figures include Karl Pribram and David Bohm discussing the holographic universe model.')
doc.save("test.docx")
print("Created test.docx")

# 4. PPTX
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Analysis: Itzhak Bentov"
subtitle.text = "Consciousness and the unified field... (Slide 1)"
prs.save("test.pptx")
print("Created test.pptx")

# 5. PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=12)
pdf.cell(200, 10, txt="McDonnell Report: Appendix A", ln=1, align="C")
pdf.multi_cell(0, 10, txt="Project Center Lane utilized various theoretical physics models, including David Bohm's implicate order and holographic brain theory by Karl Pribram, as well as the Torus model proposed by Itzhak Bentov.")
pdf.output("test.pdf")
print("Created test.pdf")
