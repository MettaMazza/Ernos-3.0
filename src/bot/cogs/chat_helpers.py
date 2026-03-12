"""
Chat Helpers - Attachment processing and reaction handling.
Extracted from ChatListener for modularity.
"""
import io
import logging
from typing import Optional

logger = logging.getLogger("ChatCog.Helpers")


class AttachmentProcessor:
    """Handles text extraction from various document formats."""
    
    @staticmethod
    async def extract_text_from_bytes(filename: str, file_bytes: bytes) -> str:
        """Extract text from document bytes. Primary method — accepts pre-read bytes.
        
        This is the preferred method since Discord attachment streams can only be
        read once. Callers should read bytes once and pass them here.
        """
        filename = filename.lower()
        file_stream = io.BytesIO(file_bytes)
        
        text = ""
        
        # PDF Handling — pdfplumber primary (handles malformed PDFs), pypdf fallback
        if filename.endswith(".pdf"):
            # Try pdfplumber first (robust against malformed cross-references)
            try:
                import pdfplumber
                with pdfplumber.open(file_stream) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                if text.strip():
                    return text
            except ImportError as e:
                logger.debug(f"Suppressed {type(e).__name__}: {e}")
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed, trying pypdf fallback: {e}")
            
            # Fallback to pypdf
            file_stream.seek(0)
            try:
                import pypdf
                reader = pypdf.PdfReader(file_stream)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            except ImportError:
                return "[Error: No PDF library installed (need pdfplumber or pypdf)]"
            except Exception as e:
                return f"[Error parsing PDF: {e}]"
        
        # DOCX Handling
        elif filename.endswith(".docx"):
            try:
                import docx
                doc = docx.Document(file_stream)
                text = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                return "[Error: python-docx not installed]"
            except Exception as e:
                return f"[Error parsing DOCX: {e}]"
                
        # PPTX Handling
        elif filename.endswith(".pptx"):
            try:
                import pptx
                prs = pptx.Presentation(file_stream)
                slide_texts = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
                text = "\n\n".join(slide_texts)
            except ImportError:
                return "[Error: python-pptx not installed]"
            except Exception as e:
                return f"[Error parsing PPTX: {e}]"
                
        # CSV Handling
        elif filename.endswith(".csv"):
            try:
                import pandas as pd
                df = pd.read_csv(file_stream)
                text = df.to_markdown(index=False)
            except ImportError:
                return "[Error: pandas or tabulate not installed]"
            except Exception as e:
                return f"[Error parsing CSV: {e}]"
                
        # XLSX Handling
        elif filename.endswith(".xlsx"):
            try:
                import pandas as pd
                # Read all sheets
                xls = pd.ExcelFile(file_stream, engine="openpyxl")
                sheet_texts = []
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    sheet_texts.append(f"=== Sheet: {sheet_name} ===\n" + df.to_markdown(index=False))
                text = "\n\n".join(sheet_texts)
            except ImportError:
                return "[Error: pandas, openpyxl, or tabulate not installed]"
            except Exception as e:
                return f"[Error parsing XLSX: {e}]"
                
        # EPUB Handling
        elif filename.endswith(".epub"):
            try:
                import ebooklib
                from ebooklib import epub
                from bs4 import BeautifulSoup as BS4
                book = epub.read_epub(file_stream)
                chapters = []
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BS4(item.get_content(), "html.parser")
                    chapter_text = soup.get_text(separator="\n", strip=True)
                    if chapter_text.strip():
                        chapters.append(chapter_text)
                text = "\n\n---\n\n".join(chapters)
            except ImportError:
                return "[Error: ebooklib or beautifulsoup4 not installed]"
            except Exception as e:
                return f"[Error parsing EPUB: {e}]"

        # ODT (OpenDocument Text) Handling
        elif filename.endswith(".odt"):
            try:
                from odf.opendocument import load as odf_load
                from odf import text as odf_text
                from odf.text import P
                doc = odf_load(file_stream)
                paragraphs = doc.getElementsByType(P)
                parts = []
                for p in paragraphs:
                    t = ""
                    for node in p.childNodes:
                        if hasattr(node, "data"):
                            t += node.data
                        elif hasattr(node, "__str__"):
                            t += str(node)
                    if t.strip():
                        parts.append(t.strip())
                text = "\n".join(parts)
            except ImportError:
                return "[Error: odfpy not installed]"
            except Exception as e:
                return f"[Error parsing ODT: {e}]"

        # ODS (OpenDocument Spreadsheet) Handling
        elif filename.endswith(".ods"):
            try:
                from odf.opendocument import load as odf_load
                from odf.table import Table, TableRow, TableCell
                from odf.text import P
                doc = odf_load(file_stream)
                sheet_texts = []
                for table in doc.getElementsByType(Table):
                    sheet_name = table.getAttribute("name") or "Sheet"
                    rows_data = []
                    for row in table.getElementsByType(TableRow):
                        cells = []
                        for cell in row.getElementsByType(TableCell):
                            cell_text = ""
                            for p in cell.getElementsByType(P):
                                for node in p.childNodes:
                                    if hasattr(node, "data"):
                                        cell_text += node.data
                            cells.append(cell_text.strip())
                        if any(cells):
                            rows_data.append(cells)
                    if rows_data:
                        # Build markdown table
                        max_cols = max(len(r) for r in rows_data)
                        rows_data = [r + [""] * (max_cols - len(r)) for r in rows_data]
                        header = "| " + " | ".join(rows_data[0]) + " |"
                        separator = "| " + " | ".join(["---"] * max_cols) + " |"
                        body = "\n".join("| " + " | ".join(r) + " |" for r in rows_data[1:])
                        sheet_texts.append(f"=== Sheet: {sheet_name} ===\n{header}\n{separator}\n{body}")
                text = "\n\n".join(sheet_texts)
            except ImportError:
                return "[Error: odfpy not installed]"
            except Exception as e:
                return f"[Error parsing ODS: {e}]"

        # ODP (OpenDocument Presentation) Handling
        elif filename.endswith(".odp"):
            try:
                from odf.opendocument import load as odf_load
                from odf.draw import Frame, Page
                from odf.text import P
                doc = odf_load(file_stream)
                slide_texts = []
                for i, page in enumerate(doc.getElementsByType(Page), 1):
                    slide_text = []
                    for p in page.getElementsByType(P):
                        t = ""
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                t += node.data
                            elif hasattr(node, "__str__"):
                                t += str(node)
                        if t.strip():
                            slide_text.append(t.strip())
                    if slide_text:
                        slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
                text = "\n\n".join(slide_texts)
            except ImportError:
                return "[Error: odfpy not installed]"
            except Exception as e:
                return f"[Error parsing ODP: {e}]"

        # Legacy DOC (Microsoft Word 97-2003) Handling
        elif filename.endswith(".doc"):
            try:
                import olefile
                ole = olefile.OleFileIO(file_stream)
                if ole.exists("WordDocument"):
                    # Extract the Word document stream
                    stream = ole.openstream("WordDocument")
                    raw = stream.read()
                    # Try to decode as text (works for most .doc files)
                    try:
                        text = raw.decode("utf-8", errors="ignore")
                    except Exception:
                        text = raw.decode("latin-1", errors="ignore")
                    # Clean up binary artifacts
                    import re
                    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
                    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
                else:
                    ole.close()
                    return "[Error: Cannot find Word content in .doc file]"
                ole.close()
            except ImportError:
                return "[Error: olefile not installed]"
            except Exception as e:
                return f"[Error parsing DOC: {e}]"

        # Legacy XLS (Microsoft Excel 97-2003) Handling
        elif filename.endswith(".xls"):
            try:
                import xlrd
                workbook = xlrd.open_workbook(file_contents=file_bytes)
                sheet_texts = []
                for sheet in workbook.sheets():
                    rows_data = []
                    for row_idx in range(sheet.nrows):
                        row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                        rows_data.append(row)
                    if rows_data:
                        max_cols = max(len(r) for r in rows_data)
                        rows_data = [r + [""] * (max_cols - len(r)) for r in rows_data]
                        header = "| " + " | ".join(rows_data[0]) + " |"
                        separator = "| " + " | ".join(["---"] * max_cols) + " |"
                        body = "\n".join("| " + " | ".join(r) + " |" for r in rows_data[1:])
                        sheet_texts.append(f"=== Sheet: {sheet.name} ===\n{header}\n{separator}\n{body}")
                text = "\n\n".join(sheet_texts)
            except ImportError:
                return "[Error: xlrd not installed]"
            except Exception as e:
                return f"[Error parsing XLS: {e}]"

        # Legacy PPT (Microsoft PowerPoint 97-2003) Handling
        elif filename.endswith(".ppt"):
            try:
                import olefile
                ole = olefile.OleFileIO(file_stream)
                if ole.exists("PowerPoint Document"):
                    stream = ole.openstream("PowerPoint Document")
                    raw = stream.read()
                    # Extract text records from the binary PPT stream
                    import re
                    # PPT stores text in UTF-16LE encoded records
                    try:
                        decoded = raw.decode("utf-16-le", errors="ignore")
                    except Exception:
                        decoded = raw.decode("latin-1", errors="ignore")
                    # Clean up binary noise
                    decoded = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', decoded)
                    lines = [line.strip() for line in decoded.splitlines() if line.strip() and len(line.strip()) > 2]
                    text = "\n".join(lines)
                else:
                    ole.close()
                    return "[Error: Cannot find PowerPoint content in .ppt file]"
                ole.close()
            except ImportError:
                return "[Error: olefile not installed]"
            except Exception as e:
                return f"[Error parsing PPT: {e}]"
                
        # Plain Text & Code Handling
        else:
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    text = file_bytes.decode("latin-1")
                except Exception:
                    return "[Error: Unknown text encoding]"
                    
        return text

    @staticmethod
    async def extract_text(attachment) -> str:
        """Extract text from a Discord attachment object.
        
        Delegates to extract_text_from_bytes. Kept for backward compatibility.
        NOTE: If you have already read the attachment bytes elsewhere, use
        extract_text_from_bytes directly to avoid the double-read issue.
        """
        file_bytes = await attachment.read()
        return await AttachmentProcessor.extract_text_from_bytes(attachment.filename, file_bytes)


class ReactionHandler:
    """Handles reaction events and social signal processing."""
    
    def __init__(self, bot):
        self.bot = bot
    
    async def process_reaction(self, payload):
        """Check for Silo Quorum and Ingest Social Signals."""
        if payload.user_id == self.bot.user.id:
            return
            
        # 1. Silo Quorum
        await self.bot.silo_manager.check_quorum(payload)
        
        # 2. Social Signal Ingestion (MRN Phase 3)
        try:
            # Check safely if cerebrum is loaded
            if hasattr(self.bot, 'cerebrum'):
                interaction_lobe = self.bot.cerebrum.lobes.get("InteractionLobe")
                if interaction_lobe:
                    social = interaction_lobe.get_ability("SocialAbility")
                    if social:
                        # Process Sentiment & Update Stats
                        sentiment = await social.process_reaction(
                            payload.user_id, 
                            str(payload.emoji), 
                            payload.message_id
                        )
                        
                        # Log to Memory (Timeline)
                        is_dm = (payload.guild_id is None)
                        if hasattr(self.bot, 'hippocampus'):
                            self.bot.hippocampus.observe_reaction(
                                user_id=str(payload.user_id),
                                emoji=str(payload.emoji),
                                sentiment=sentiment,
                                message_id=payload.message_id,
                                channel_id=payload.channel_id,
                                is_dm=is_dm
                            )
        except Exception as e:
            logger.error(f"Failed to process reaction signal: {e}")
