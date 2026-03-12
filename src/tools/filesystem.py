import os
import logging
from .registry import ToolRegistry
from src.privacy.guard import validate_path_scope
from src.privacy.scopes import PrivacyScope

logger = logging.getLogger("Tools.Filesystem")

@ToolRegistry.register(name="read_file_page", description="Read a specific section of a large file.")
def read_file_page(path: str, start_line: int = 1, limit: int = 5000, request_scope: str = "PUBLIC", user_id: str = None) -> str:
    """Reads lines from a file with scope and user ownership validation."""
    # SCOPE CHECK: Validate path access FIRST
    try:
        scope = PrivacyScope[request_scope.upper()]
    except Exception:
        scope = PrivacyScope.PUBLIC
        
    if not validate_path_scope(path, scope, user_id=user_id):
        logger.warning(f"Scope violation: {request_scope} (user={user_id}) tried to read {path}")
        return (
            f"🔒 Access Denied: Your scope ({request_scope}) cannot access this path.\n"
            f"[CRITICAL]: You MUST report this access limitation honestly. "
            f"Do NOT substitute web search results, fabricated content, or paraphrased "
            f"data in place of this denied file. State clearly that this content is not "
            f"accessible in the current scope."
        )

    if not os.path.exists(path):
        return f"Error: File not found at {path}"
    
    try:
        lines = []
        ext = path.lower()
        if ext.endswith(".docx"):
            try:
                import docx
                doc = docx.Document(path)
                lines = [p.text + '\n' for p in doc.paragraphs if p.text.strip()]
            except ImportError:
                return f"Error: python-docx not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .docx file: {e}"
        elif ext.endswith(".pptx"):
            try:
                import pptx
                prs = pptx.Presentation(path)
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        lines.append(f"--- Slide {i} ---\n")
                        lines.extend([t + '\n' for t in slide_text])
                        lines.append("\n")
            except ImportError:
                return f"Error: python-pptx not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .pptx file: {e}"
        elif ext.endswith(".csv"):
            try:
                import pandas as pd
                df = pd.read_csv(path)
                lines = df.to_markdown(index=False).splitlines(keepends=True)
                lines = [line if line.endswith('\n') else line + '\n' for line in lines]
            except ImportError:
                return f"Error: pandas or tabulate not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .csv file: {e}"
        elif ext.endswith(".xlsx"):
            try:
                import pandas as pd
                xls = pd.ExcelFile(path, engine="openpyxl")
                for sheet_name in xls.sheet_names:
                    lines.append(f"=== Sheet: {sheet_name} ===\n")
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    sheet_lines = df.to_markdown(index=False).splitlines(keepends=True)
                    sheet_lines = [line if line.endswith('\n') else line + '\n' for line in sheet_lines]
                    lines.extend(sheet_lines)
                    lines.append("\n")
            except ImportError:
                return f"Error: pandas, openpyxl, or tabulate not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .xlsx file: {e}"
        elif ext.endswith(".pdf"):
            try:
                import fitz
                doc = fitz.open(path)
                for i, page in enumerate(doc):
                    lines.append(f"--- Page {i+1} ---\n")
                    text = page.get_text()
                    lines.extend([line + '\n' for line in text.splitlines() if line.strip()])
                    lines.append("\n")
            except ImportError:
                return f"Error: PyMuPDF (fitz) not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .pdf file: {e}"
        elif ext.endswith(".epub"):
            try:
                import ebooklib
                from ebooklib import epub
                from bs4 import BeautifulSoup as BS4
                book = epub.read_epub(path)
                chapters = []
                for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                    soup = BS4(item.get_content(), "html.parser")
                    chapter_text = soup.get_text(separator="\n", strip=True)
                    if chapter_text.strip():
                        chapters.append(chapter_text)
                lines = ("\n\n---\n\n".join(chapters)).splitlines(keepends=True)
                lines = [line if line.endswith('\n') else line + '\n' for line in lines]
            except ImportError:
                return f"Error: ebooklib or beautifulsoup4 not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .epub file: {e}"
        elif ext.endswith(".odt"):
            try:
                from odf.opendocument import load as odf_load
                from odf.text import P
                doc = odf_load(path)
                paragraphs = doc.getElementsByType(P)
                parts = []
                for p in paragraphs:
                    t = ""
                    for node in p.childNodes:
                        if hasattr(node, "data"):
                            t += node.data
                        elif hasattr(node, "__str__"):
                            t += str(node)
                    if t.strip():
                        parts.append(t.strip())
                lines = [p + '\n' for p in parts]
            except ImportError:
                return f"Error: odfpy not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .odt file: {e}"
        elif ext.endswith(".ods"):
            try:
                from odf.opendocument import load as odf_load
                from odf.table import Table, TableRow, TableCell
                from odf.text import P
                doc = odf_load(path)
                all_lines = []
                for table in doc.getElementsByType(Table):
                    sheet_name = table.getAttribute("name") or "Sheet"
                    all_lines.append(f"=== Sheet: {sheet_name} ===\n")
                    rows_data = []
                    for row in table.getElementsByType(TableRow):
                        cells = []
                        for cell in row.getElementsByType(TableCell):
                            cell_text = ""
                            for p in cell.getElementsByType(P):
                                for node in p.childNodes:
                                    if hasattr(node, "data"):
                                        cell_text += node.data
                            cells.append(cell_text.strip())
                        if any(cells):
                            rows_data.append(cells)
                    if rows_data:
                        max_cols = max(len(r) for r in rows_data)
                        rows_data = [r + [""] * (max_cols - len(r)) for r in rows_data]
                        all_lines.append("| " + " | ".join(rows_data[0]) + " |\n")
                        all_lines.append("| " + " | ".join(["---"] * max_cols) + " |\n")
                        for r in rows_data[1:]:
                            all_lines.append("| " + " | ".join(r) + " |\n")
                    all_lines.append("\n")
                lines = all_lines
            except ImportError:
                return f"Error: odfpy not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .ods file: {e}"
        elif ext.endswith(".odp"):
            try:
                from odf.opendocument import load as odf_load
                from odf.draw import Page
                from odf.text import P
                doc = odf_load(path)
                all_lines = []
                for i, page in enumerate(doc.getElementsByType(Page), 1):
                    slide_text = []
                    for p in page.getElementsByType(P):
                        t = ""
                        for node in p.childNodes:
                            if hasattr(node, "data"):
                                t += node.data
                            elif hasattr(node, "__str__"):
                                t += str(node)
                        if t.strip():
                            slide_text.append(t.strip())
                    if slide_text:
                        all_lines.append(f"--- Slide {i} ---\n")
                        all_lines.extend([t + '\n' for t in slide_text])
                        all_lines.append("\n")
                lines = all_lines
            except ImportError:
                return f"Error: odfpy not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .odp file: {e}"
        elif ext.endswith(".doc"):
            try:
                import olefile
                ole = olefile.OleFileIO(path)
                if ole.exists("WordDocument"):
                    stream = ole.openstream("WordDocument")
                    raw = stream.read()
                    try:
                        content = raw.decode("utf-8", errors="ignore")
                    except Exception:
                        content = raw.decode("latin-1", errors="ignore")
                    import re
                    content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', content)
                    lines = [line.strip() + '\n' for line in content.splitlines() if line.strip()]
                else:
                    ole.close()
                    return f"Error: Cannot find Word content in .doc file: {path}"
                ole.close()
            except ImportError:
                return f"Error: olefile not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .doc file: {e}"
        elif ext.endswith(".xls"):
            try:
                import xlrd
                workbook = xlrd.open_workbook(path)
                all_lines = []
                for sheet in workbook.sheets():
                    all_lines.append(f"=== Sheet: {sheet.name} ===\n")
                    rows_data = []
                    for row_idx in range(sheet.nrows):
                        row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                        rows_data.append(row)
                    if rows_data:
                        max_cols = max(len(r) for r in rows_data)
                        rows_data = [r + [""] * (max_cols - len(r)) for r in rows_data]
                        all_lines.append("| " + " | ".join(rows_data[0]) + " |\n")
                        all_lines.append("| " + " | ".join(["---"] * max_cols) + " |\n")
                        for r in rows_data[1:]:
                            all_lines.append("| " + " | ".join(r) + " |\n")
                    all_lines.append("\n")
                lines = all_lines
            except ImportError:
                return f"Error: xlrd not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .xls file: {e}"
        elif ext.endswith(".ppt"):
            try:
                import olefile
                ole = olefile.OleFileIO(path)
                if ole.exists("PowerPoint Document"):
                    stream = ole.openstream("PowerPoint Document")
                    raw = stream.read()
                    import re
                    try:
                        decoded = raw.decode("utf-16-le", errors="ignore")
                    except Exception:
                        decoded = raw.decode("latin-1", errors="ignore")
                    decoded = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', decoded)
                    lines = [line.strip() + '\n' for line in decoded.splitlines() if line.strip() and len(line.strip()) > 2]
                else:
                    ole.close()
                    return f"Error: Cannot find PowerPoint content in .ppt file: {path}"
                ole.close()
            except ImportError:
                return f"Error: olefile not installed. Could not read {path}"
            except Exception as e:
                return f"Error reading .ppt file: {e}"
        else:
            with open(path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
            
        total_lines = len(lines)
        start_index = max(0, start_line - 1)
        end_index = min(total_lines, start_index + limit)
        
        content = "".join(lines[start_index:end_index])
        
        # ─── Reading Progress Metadata ──────────────────────
        pct = int((end_index / total_lines) * 100) if total_lines > 0 else 100
        remaining = total_lines - end_index
        
        header = f"File: {path}\nLines: {start_index+1}-{end_index}/{total_lines} ({pct}% complete"
        if remaining > 0:
            header += f", {remaining} lines remaining)\n[BOOKMARK: Continue with read_file(path='{path}', start_line={end_index + 1})]\n[READING INCOMPLETE — you MUST continue reading before responding]"
        else:
            header += ")\n[DOCUMENT COMPLETE]"
        
        return f"{header}\n\n{content}"
    except Exception as e:
        return f"Error reading file: {e}"

@ToolRegistry.register(name="search_codebase", description="Search for string in files.")
def search_codebase(query: str, path: str = "./src", request_scope: str = "PUBLIC", user_id: str = None) -> str:
    """Simple grep-like search with scope and user ownership validation."""
    # SCOPE CHECK: Block searching memory directories from PUBLIC
    try:
        scope = PrivacyScope[request_scope.upper()]
    except Exception:
        scope = PrivacyScope.PUBLIC
        
    if not validate_path_scope(path, scope, user_id=user_id):
        logger.warning(f"Scope violation: {request_scope} (user={user_id}) tried to search {path}")
        return (
            f"🔒 Access Denied: Your scope ({request_scope}) cannot search this path.\n"
            f"[CRITICAL]: You MUST report this access limitation honestly. "
            f"Do NOT substitute web search results, fabricated content, or paraphrased "
            f"data in place of this denied file. State clearly that this content is not "
            f"accessible in the current scope."
        )
    
    results = []
    try:
        for root, dirs, files in os.walk(path):
            # Skip memory directories for non-CORE scope
            if scope != PrivacyScope.CORE_PRIVATE:
                dirs[:] = [d for d in dirs if d not in ['core']]
                
            for file in files:
                if file.endswith((".py", ".md", ".txt")):
                    file_path = os.path.join(root, file)
                    
                    # Additional path check for each file
                    if not validate_path_scope(file_path, scope, user_id=user_id):
                        continue
                        
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        for i, line in enumerate(f, 1):
                            if query in line:
                                results.append(f"{file_path}:{i}: {line.strip()}")
                                if len(results) >= 20:
                                    return "\n".join(results) + "\n... (truncated)"
        return "\n".join(results) if results else "No matches found."
    except Exception as e:
        return f"Search Error: {e}"

@ToolRegistry.register(name="read_file", description="Read a file. Alias for read_file_page.")
def read_file(path: str, start_line: int = 1, limit: int = 5000, request_scope: str = "PUBLIC", user_id: str = None) -> str:
    """Alias for read_file_page with scope validation."""
    return read_file_page(path, start_line, limit, request_scope, user_id=user_id)

@ToolRegistry.register(name="list_files", description="List files in a directory.")
def list_files(path: str = ".", request_scope: str = "PUBLIC", user_id: str = None) -> str:
    """List files in a directory with scope and user ownership validation."""
    # SCOPE CHECK
    try:
        scope = PrivacyScope[request_scope.upper()]
    except Exception:
        scope = PrivacyScope.PUBLIC
        
    if not validate_path_scope(path, scope, user_id=user_id):
        logger.warning(f"Scope violation: {request_scope} (user={user_id}) tried to list {path}")
        return (
            f"🔒 Access Denied: Your scope ({request_scope}) cannot list this path.\n"
            f"[CRITICAL]: You MUST report this access limitation honestly. "
            f"Do NOT substitute web search results, fabricated content, or paraphrased "
            f"data in place of this denied file. State clearly that this content is not "
            f"accessible in the current scope."
        )
    
    try:
        if not os.path.exists(path):
            return f"Error: Path not found: {path}"
        if os.path.isfile(path):
            return f"{path} (file)"
        
        entries = []
        for entry in os.listdir(path):
            full_path = os.path.join(path, entry)
            
            # Filter out inaccessible directories
            if not validate_path_scope(full_path, scope, user_id=user_id):
                continue
                
            if os.path.isdir(full_path):
                entries.append(f"[DIR] {entry}/")
            else:
                entries.append(f"      {entry}")
        return f"Contents of {path}:\n" + "\n".join(sorted(entries))
    except Exception as e:
        return f"Error listing files: {e}"
